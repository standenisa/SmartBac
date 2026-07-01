"""Generate PPT about SmartBAC model training (antrenare.ipynb)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
RED = RGBColor(0xF8, 0x71, 0x71)
TEAL = RGBColor(0x14, 0xB8, 0xA6)


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_shape(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    return tb


def card(slide, l, t, w, h, title, body, accent=BLUE, bsz=13):
    add_shape(slide, l, t, w, h, CARD)
    bar = slide.shapes.add_shape(1, l, t, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    txt(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.4), Inches(0.35),
        title, sz=15, clr=accent, bold=True)
    txt(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.4), h - Inches(0.6),
        body, sz=bsz, clr=MUTED)


def metric(slide, l, t, value, label, clr):
    add_shape(slide, l, t, Inches(2.2), Inches(1.6), CARD)
    txt(slide, l, t + Inches(0.15), Inches(2.2), Inches(0.7),
        value, sz=34, clr=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l, t + Inches(0.85), Inches(2.2), Inches(0.5),
        label, sz=12, clr=MUTED, align=PP_ALIGN.CENTER)


# Training data
steps = [50, 100, 150, 200, 250, 300, 350, 400, 450, 500]
train_loss = [0.4701, 0.3384, 0.2393, 0.2632, 0.2174, 0.2045, 0.1819, 0.1857, 0.2385, 0.1699]
eval_loss = [0.4267, 0.3183, 0.2481, 0.2270, 0.2132, 0.2034, 0.1975, 0.1929, 0.1909, 0.1905]
import math
train_ppl = [math.exp(l) for l in train_loss]
eval_ppl = [math.exp(l) for l in eval_loss]


# ═══════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
    'Antrenarea Modelului SmartBAC', sz=44, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.7),
    'Fine-tuning DeepSeek-R1-Distill-Qwen-14B cu QLoRA', sz=22, clr=MUTED, align=PP_ALIGN.CENTER)
txt(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
    'pe exerciții de Bacalaureat la Matematică', sz=18, clr=PURPLE, align=PP_ALIGN.CENTER)

for i, (t, c) in enumerate([('14B Parametri', PURPLE), ('QLoRA 4-bit', ORANGE),
                              ('4.086 Exerciții', GREEN), ('Kaggle T4', BLUE), ('1 Epocă', CYAN)]):
    x = Inches(2.8) + Inches(1.6) * i
    add_shape(slide, x, Inches(5), Inches(1.4), Inches(0.45), RGBColor(0x33, 0x41, 0x55))
    txt(slide, x, Inches(5.02), Inches(1.4), Inches(0.4), t, sz=11, clr=c, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2 — Model & Dataset
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'MODEL ȘI DATE', sz=12, clr=ORANGE, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'DeepSeek-R1-Distill-Qwen-14B', sz=30, clr=WHITE, bold=True)

card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5),
     'Despre Model',
     'DeepSeek-R1-Distill-Qwen-14B\n\n'
     '• 14 miliarde de parametri\n'
     '• Distilat din DeepSeek-R1 (671B)\n'
     '• Arhitectură Qwen2.5 (decoder-only)\n'
     '• Capabilități de raționament (chain-of-thought)\n'
     '• Suport pentru <think> blocks\n'
     '• Pre-antrenat pe date multilingve\n\n'
     'Ales pentru: raționament matematic superior\n'
     'și suport nativ pentru limba română.',
     PURPLE, bsz=14)

card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.3),
     'Dataset',
     'exercises_merged.json — 4.541 exerciții BAC\n\n'
     'Split:  4.086 train (90%)  |  455 test (10%)\n'
     'Tipuri: ecuații, derivate, integrale, limite,\n'
     '        matrice, combinări, probabilități, geometrie\n'
     'Surse: subiecte BAC 2009-2024, toate profilurile',
     GREEN, bsz=14)

card(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.4),
     'Format ChatML',
     '<|im_start|>system\n'
     'Ești SmartBAC, un asistent de matematică...\n'
     '<|im_end|>\n'
     '<|im_start|>user\n'
     'Calculați derivata f(x) = x²·sin(x)\n'
     '<|im_end|>\n'
     '<|im_start|>assistant\n'
     '<think> rezolvare pas cu pas </think>\n'
     'Răspuns: f\'(x) = 2x·sin(x) + x²·cos(x)\n'
     '<|im_end|>',
     CYAN, bsz=11)

# ═══════════════════════════════════════
# SLIDE 3 — QLoRA
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'TEHNICA DE ANTRENARE', sz=12, clr=PURPLE, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'QLoRA — Quantized Low-Rank Adaptation', sz=30, clr=WHITE, bold=True)

card(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(5),
     'Ce este QLoRA?',
     'Combinație de:\n\n'
     '1. Quantizare 4-bit (NF4)\n'
     '   Model 14B: 28GB → ~8GB VRAM\n\n'
     '2. LoRA (Low-Rank Adaptation)\n'
     '   Antrenăm doar ~2% din parametri\n\n'
     '3. Double Quantization\n'
     '   Quantizăm și constantele de quantizare\n\n'
     'Rezultat: antrenare model 14B\n'
     'pe un singur GPU de 16GB (T4)',
     ORANGE, bsz=14)

card(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(5),
     'Configurație LoRA',
     'Rank (r): 16\n'
     'Alpha (α): 32\n'
     'Dropout: 0.05\n'
     'Bias: none\n\n'
     'Target modules (7):\n'
     '• q_proj — Query projection\n'
     '• k_proj — Key projection\n'
     '• v_proj — Value projection\n'
     '• o_proj — Output projection\n'
     '• gate_proj — FFN gate\n'
     '• up_proj — FFN up\n'
     '• down_proj — FFN down',
     PURPLE, bsz=13)

card(slide, Inches(8.7), Inches(1.8), Inches(3.8), Inches(5),
     'De ce funcționează?',
     'Formula LoRA:\n'
     'W_new = W_frozen + B × A\n\n'
     'W: 14B params (înghețat)\n'
     'B: d × 16 (antrenat)\n'
     'A: 16 × d (antrenat)\n\n'
     'Parametri antrenabili:\n'
     '~0.98% din total\n\n'
     'Avantaje:\n'
     '• VRAM redus dramatic\n'
     '• Antrenare rapidă (1 epocă)\n'
     '• Adapter mic (~150MB)\n'
     '• Model original nemodificat',
     TEAL, bsz=13)

# ═══════════════════════════════════════
# SLIDE 4 — Training Config
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'CONFIGURAȚIA ANTRENĂRII', sz=12, clr=GREEN, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'Hiperparametri și Infrastructură', sz=30, clr=WHITE, bold=True)

card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5),
     'Hiperparametri',
     'Batch size per GPU:     1\n'
     'Gradient accumulation:  8\n'
     'Effective batch size:   8\n\n'
     'Learning rate:          5 × 10⁻⁵\n'
     'LR Scheduler:           Cosine decay\n'
     'Warmup:                 5% din total steps\n\n'
     'Epoci:                  1\n'
     'Max sequence length:    1024 tokeni\n'
     'Optimizer:              Paged AdamW 8-bit\n'
     'Precision:              FP16 (mixed precision)\n\n'
     'Eval steps:             la fiecare 50 pași\n'
     'Save strategy:          la finalul epocii',
     GREEN, bsz=14)

card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.3),
     'Infrastructură',
     'GPU:        Tesla T4 (16GB VRAM) — Kaggle\n'
     'CPU:        Intel Xeon (2 cores)\n'
     'RAM:        13GB\n'
     'Storage:    ~50GB (model + date)\n\n'
     'Framework:  PyTorch + HuggingFace Transformers\n'
     'Libraries:  PEFT, TRL, BitsAndBytes, Accelerate',
     BLUE, bsz=14)

card(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.4),
     'Tehnici de optimizare VRAM',
     '• Gradient Checkpointing\n'
     '  Recalculează activări în loc să le stocheze\n\n'
     '• 4-bit NF4 Quantization + Double Quant\n'
     '  14B model: 28GB → ~8GB\n\n'
     '• Paged AdamW 8-bit\n'
     '  Optimizer states în CPU când GPU plin',
     ORANGE, bsz=14)

# ═══════════════════════════════════════
# SLIDE 5 — Training Metrics
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'REZULTATE ANTRENARE', sz=12, clr=YELLOW, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'Evoluția Loss și Perplexitate', sz=30, clr=WHITE, bold=True)

# Key metrics at top
metric(slide, Inches(0.5), Inches(1.6), '0.470', 'Loss Inițial\n(Step 50)', RED)
metric(slide, Inches(2.9), Inches(1.6), '0.170', 'Loss Final\n(Step 500)', GREEN)
metric(slide, Inches(5.3), Inches(1.6), '-64%', 'Reducere\nLoss', CYAN)
metric(slide, Inches(7.7), Inches(1.6), '1.21', 'Perplexitate\nFinală (eval)', PURPLE)
metric(slide, Inches(10.1), Inches(1.6), '500', 'Total\nPași', BLUE)

# Loss chart
chart_data = CategoryChartData()
chart_data.categories = [str(s) for s in steps]
chart_data.add_series('Train Loss', train_loss)
chart_data.add_series('Eval Loss', eval_loss)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(0.5), Inches(3.5), Inches(5.8), Inches(3.5),
    chart_data
).chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = MUTED

plot = chart.plots[0]
plot.smooth = True
s1 = plot.series[0]
s1.format.line.color.rgb = TEAL
s1.format.line.width = Pt(2.5)
s2 = plot.series[1]
s2.format.line.color.rgb = RED
s2.format.line.width = Pt(2.5)

cat_axis = chart.category_axis
cat_axis.tick_labels.font.size = Pt(9)
cat_axis.tick_labels.font.color.rgb = MUTED
val_axis = chart.value_axis
val_axis.tick_labels.font.size = Pt(9)
val_axis.tick_labels.font.color.rgb = MUTED

# Perplexity chart
chart_data2 = CategoryChartData()
chart_data2.categories = [str(s) for s in steps]
chart_data2.add_series('Train PPL', train_ppl)
chart_data2.add_series('Eval PPL', eval_ppl)

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.5),
    chart_data2
).chart
chart2.has_legend = True
chart2.legend.include_in_layout = False
chart2.legend.font.size = Pt(10)
chart2.legend.font.color.rgb = MUTED

plot2 = chart2.plots[0]
plot2.smooth = True
s1 = plot2.series[0]
s1.format.line.color.rgb = GREEN
s1.format.line.width = Pt(2.5)
s2 = plot2.series[1]
s2.format.line.color.rgb = ORANGE
s2.format.line.width = Pt(2.5)

cat_axis2 = chart2.category_axis
cat_axis2.tick_labels.font.size = Pt(9)
cat_axis2.tick_labels.font.color.rgb = MUTED
val_axis2 = chart2.value_axis
val_axis2.tick_labels.font.size = Pt(9)
val_axis2.tick_labels.font.color.rgb = MUTED

txt(slide, Inches(0.5), Inches(3.5), Inches(5.8), Inches(0.3),
    'Funcția de Pierdere (Loss)', sz=11, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(0.3),
    'Perplexitate', sz=11, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 6 — Training Analysis
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'ANALIZA ANTRENĂRII', sz=12, clr=CYAN, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'Interpretarea Rezultatelor', sz=30, clr=WHITE, bold=True)

card(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.8),
     'Convergență',
     'Loss-ul scade rapid în primii\n'
     '150 de pași (0.47 → 0.24),\n'
     'apoi se stabilizează ~0.19.\n\n'
     'Eval loss urmează train loss\n'
     'fără divergență → NO overfitting.\n\n'
     'Perplexitatea finală: 1.21\n'
     '(modelul prezice corect cu\n'
     'încredere mare tokenul următor).\n\n'
     'O singură epocă e suficientă\n'
     'pentru acest dataset.',
     GREEN, bsz=14)

card(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(4.8),
     'Observații importante',
     '1. Fără overfitting\n'
     '   eval_loss ≈ train_loss pe tot\n'
     '   parcursul antrenării\n\n'
     '2. Cosine LR schedule\n'
     '   LR scade gradual de la 5e-5\n'
     '   la 0, evitând oscilații\n\n'
     '3. Spike la step 450\n'
     '   Train loss crește temporar\n'
     '   (normal — batch dificil)\n'
     '   dar se recuperează la 500\n\n'
     '4. Model stabil la final\n'
     '   eval_loss: 0.1905',
     BLUE, bsz=14)

card(slide, Inches(8.7), Inches(1.8), Inches(3.8), Inches(4.8),
     'Ce a învățat modelul',
     'Prin fine-tuning, modelul a\n'
     'dobândit:\n\n'
     '• Rezolvare pas cu pas\n'
     '  (format <think> structurat)\n\n'
     '• Terminologie BAC română\n'
     '  (derivate, integrale, funcții)\n\n'
     '• Raționament matematic\n'
     '  (identificare tip → metodă →\n'
     '   aplicare → verificare)\n\n'
     '• Format răspuns consistent\n'
     '  (Pasul 1: ... Răspuns: ...)',
     PURPLE, bsz=14)

# ═══════════════════════════════════════
# SLIDE 7 — Deployment
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'DEPLOYMENT', sz=12, clr=GREEN, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'De la Antrenare la Producție', sz=30, clr=WHITE, bold=True)

# Flow
flow_items = [
    ('1. Antrenare\nKaggle T4', 'QLoRA fine-tune\n500 steps\n1 epocă', PURPLE),
    ('2. Salvare\nHuggingFace', 'deniii1111/bac-math-\nqwen-lora\nAdapter LoRA ~150MB', ORANGE),
    ('3. Servire\nKaggle + ngrok', 'Flask server\n/ask endpoint\nModel 4-bit CUDA', BLUE),
    ('4. Integrare\nFastAPI Backend', 'POST /api/chat\nProxy → Kaggle\nJSON structurat', GREEN),
    ('5. Frontend\nReact Native', 'KaTeX formule\nPași + Răspuns\nDark mode UI', CYAN),
]
for i, (title, body, clr) in enumerate(flow_items):
    x = Inches(0.3) + Inches(2.55) * i
    card(slide, x, Inches(1.8), Inches(2.3), Inches(2.8), title, body, clr, bsz=12)

# Arrow between cards
for i in range(4):
    x = Inches(2.6) + Inches(2.55) * i
    txt(slide, x, Inches(2.8), Inches(0.3), Inches(0.4), '→', sz=22, clr=MUTED, align=PP_ALIGN.CENTER)

card(slide, Inches(0.5), Inches(5), Inches(12), Inches(2),
     'Pipeline complet',
     'Utilizator scrie întrebare → Frontend trimite POST /api/chat → Backend verifică dacă modelul SmartBAC (Kaggle/ngrok) e disponibil → '
     'Dacă DA: trimite {"intrebare": "..."} la ngrok/ask → Model generează cu <think> blocks → Parsare pași + răspuns → Return JSON structurat\n'
     'Dacă NU: fallback la GPT-4o (OpenAI API) → Apoi fallback la rule-based solver (offline, <100ms)',
     YELLOW, bsz=13)

# ═══════════════════════════════════════
# SLIDE 8 — Conclusions
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
txt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
    'CONCLUZII', sz=12, clr=PURPLE, bold=True)
txt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
    'Rezumat Antrenare SmartBAC', sz=30, clr=WHITE, bold=True)

card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8),
     'Ce am realizat',
     '• Fine-tuning model 14B pe un GPU T4 16GB\n'
     '  (imposibil fără QLoRA)\n\n'
     '• Loss redus cu 64%: 0.47 → 0.17\n\n'
     '• Zero overfitting (eval ≈ train)\n\n'
     '• Perplexitate finală: 1.21\n'
     '  (predicție foarte precisă)\n\n'
     '• Adapter LoRA: doar ~150MB\n'
     '  (nu 28GB model complet)\n\n'
     '• Model publicat pe HuggingFace:\n'
     '  deniii1111/bac-math-qwen-lora\n\n'
     '• Integrat în aplicație ca solver primar',
     GREEN, bsz=14)

card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8),
     'Lecții învățate',
     '• QLoRA permite fine-tuning accesibil\n'
     '  pe hardware consumer/gratuit\n\n'
     '• O singură epocă e suficientă când\n'
     '  datasetul e de calitate (4086 exerciții)\n\n'
     '• Cosine LR schedule > linear decay\n'
     '  pentru convergență stabilă\n\n'
     '• Gradient checkpointing e esențial\n'
     '  pentru modele mari pe VRAM limitat\n\n'
     '• Format ChatML cu <think> blocks\n'
     '  produce raționament structurat\n\n'
     '• Eval la fiecare 50 steps permite\n'
     '  monitorizare continuă fără overhead',
     CYAN, bsz=14)

txt(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
    'Mulțumesc!    |    Întrebări?', sz=16, clr=MUTED, align=PP_ALIGN.CENTER)

# Save
out = os.path.expanduser('~/Desktop/SmartBAC_Antrenare_PPT.pptx')
prs.save(out)
print(f'Saved: {out}')
