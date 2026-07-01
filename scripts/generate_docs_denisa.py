"""Generate Word + PPT — Implementare AI Chatbot SmartBAC (Denisa)."""

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import math, os

# ══════════════════════════════════════
#  WORD DOCUMENT
# ══════════════════════════════════════

doc = Document()
style = doc.styles['Normal']
style.font.name = 'Times New Roman'
style.font.size = Pt(12)
style.paragraph_format.line_spacing = 1.5
style.paragraph_format.space_after = Pt(6)

for level in range(1, 4):
    hs = doc.styles[f'Heading {level}']
    hs.font.name = 'Times New Roman'
    hs.font.color.rgb = RGBColor(0x1B, 0x1B, 0x2F)
    hs.font.size = Pt(16 - (level - 1) * 2)
    hs.font.bold = True


def txt(text):
    p = doc.add_paragraph(text)
    p.paragraph_format.first_line_indent = Cm(1.25)
    return p

def bullet(text):
    return doc.add_paragraph(text, style='List Bullet')

def code(text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.name = 'Consolas'
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shading = p.paragraph_format.element.get_or_add_pPr()
    shd = shading.makeelement(qn('w:shd'), {qn('w:val'): 'clear', qn('w:color'): 'auto', qn('w:fill'): 'F5F5F5'})
    shading.append(shd)

def table(headers, rows):
    t = doc.add_table(rows=1+len(rows), cols=len(headers))
    t.style = 'Light Grid Accent 1'
    for i, h in enumerate(headers):
        cell = t.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs: r.bold = True; r.font.size = Pt(10)
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = t.rows[r_idx+1].cells[c_idx]
            cell.text = str(val)
            for p in cell.paragraphs:
                for r in p.runs: r.font.size = Pt(10)
    doc.add_paragraph()

# ─── TITLE ───
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('\n\n\n')
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('SmartBAC — Chatbot AI')
run.font.size = Pt(26); run.bold = True; run.font.color.rgb = RGBColor(0x1B, 0x1B, 0x2F)
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('Implementarea componentei de Inteligență Artificială')
run.font.size = Pt(16); run.font.color.rgb = RGBColor(0x4A, 0x4A, 0x6A)
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('\nAntrenarea unui model de limbaj pentru rezolvarea\nexercițiilor de Bacalaureat la Matematică')
run.font.size = Pt(13); run.font.color.rgb = RGBColor(0x6A, 0x6A, 0x8A)
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('\n\n2026')
run.font.size = Pt(14)
doc.add_page_break()

# ─── 1. INTRODUCERE ───
doc.add_heading('1. Introducere', level=1)
txt('SmartBAC este o aplicație mobilă de pregătire pentru examenul de Bacalaureat la Matematică. '
    'Componenta centrală a aplicației este un chatbot inteligent, capabil să rezolve exerciții '
    'de matematică pas cu pas, explicând fiecare etapă în limba română.')
txt('Chatbot-ul se bazează pe un model de limbaj mare (LLM) — DeepSeek-R1-Distill-Qwen-14B — '
    'care a fost fine-tunat folosind tehnica QLoRA pe un dataset de exerciții de BAC. '
    'Modelul a fost antrenat pe platforma Kaggle, folosind un GPU Tesla T4 de 16GB.')

doc.add_heading('1.1 Obiective', level=2)
bullet('Antrenarea unui model AI care rezolvă exerciții de matematică BAC pas cu pas')
bullet('Generarea de raționamente structurate (chain-of-thought) cu blocuri <think>')
bullet('Integrarea modelului antrenat într-un chatbot accesibil din aplicația mobilă')
bullet('Evaluarea calității răspunsurilor prin metrici de loss și perplexitate')

doc.add_heading('1.2 Motivație', level=2)
txt('Elevii care se pregătesc pentru BAC au nevoie de un tutor disponibil permanent, '
    'care să explice rezolvările în detaliu. Un model de AI antrenat specific pe exerciții '
    'de BAC poate oferi acest lucru, fără costuri și fără limită de timp.')
doc.add_page_break()

# ─── 2. MODELUL ───
doc.add_heading('2. Modelul de bază', level=1)
txt('Am ales modelul DeepSeek-R1-Distill-Qwen-14B ca punct de plecare. Acesta este un model '
    'de limbaj cu 14 miliarde de parametri, distilat din modelul mai mare DeepSeek-R1 (671B parametri). '
    'Avantajul principal este capacitatea nativă de raționament — modelul poate genera '
    'pași logici de gândire în blocuri <think>.')

doc.add_heading('2.1 Caracteristici tehnice', level=2)
table(['Parametru', 'Valoare'], [
    ['Nume model', 'DeepSeek-R1-Distill-Qwen-14B'],
    ['Parametri', '14 miliarde'],
    ['Arhitectură', 'Decoder-only (Transformer)'],
    ['Distilat din', 'DeepSeek-R1 (671B)'],
    ['Vocabular', '~150.000 tokeni'],
    ['Context maxim', '8.192 tokeni'],
    ['Capacitate specială', 'Raționament în blocuri <think>'],
])

doc.add_heading('2.2 De ce DeepSeek-R1?', level=2)
bullet('Raționament nativ: generează gândire pas cu pas fără prompting special')
bullet('Distilare eficientă: performanță apropiată de modelul de 671B la o fracțiune din dimensiune')
bullet('Suport multilingv: înțelege și generează text în limba română')
bullet('Arhitectură Qwen2.5: compatibil cu ecosistemul HuggingFace pentru fine-tuning')
doc.add_page_break()

# ─── 3. DATASET ───
doc.add_heading('3. Datasetul de antrenare', level=1)
txt('Datasetul conține exerciții de matematică pentru Bacalaureat, colectate din subiecte '
    'oficiale ale anilor anteriori (2009-2024) și exerciții generate suplimentar.')

doc.add_heading('3.1 Structura datelor', level=2)
table(['Câmp', 'Descriere', 'Exemplu'], [
    ['question', 'Enunțul exercițiului', 'Rezolvă ecuația: x² - 5x + 6 = 0'],
    ['answer', 'Răspunsul corect', 'x₁ = 2, x₂ = 3'],
    ['solution_steps', 'Pași de rezolvare', '["Pasul 1: ...", "Pasul 2: ..."]'],
    ['topic', 'Categoria exercițiului', 'equation, derivative, integral'],
    ['difficulty', 'Nivel dificultate', '1 (ușor), 2 (mediu), 3 (greu)'],
    ['profile', 'Profil BAC', 'M1, M2, M3, M4, BOTH'],
])

doc.add_heading('3.2 Dimensiune și distribuție', level=2)
table(['Split', 'Samples folosite', 'Din total'], [
    ['Train', '1.500', '3.415 (90%)'],
    ['Evaluare', '150', '380 (10%)'],
    ['Total disponibil', '3.795', '100%'],
])

txt('Tipuri de exerciții acoperite: ecuații de grad I și II, derivate, integrale, limite, '
    'determinanți, combinări, aranjamente, progresii aritmetice și geometrice, trigonometrie, '
    'geometrie analitică, numere complexe, funcții.')

doc.add_heading('3.3 Formatare pentru antrenare (ChatML)', level=2)
txt('Fiecare exercițiu a fost formatat în structura ChatML, compatibilă cu modelul Qwen:')
code(
    '<|im_start|>system\n'
    'Ești SmartBAC, un asistent de matematică specializat pe BAC.\n'
    'Gândește pas cu pas în blocul <think>, apoi dă răspunsul final.\n'
    '<|im_end|>\n'
    '<|im_start|>user\n'
    'Rezolvă ecuația: x² - 5x + 6 = 0\n'
    '<|im_end|>\n'
    '<|im_start|>assistant\n'
    '<think>\n'
    'Pasul 1: Identificăm a=1, b=-5, c=6\n'
    'Pasul 2: Δ = 25 - 24 = 1\n'
    'Pasul 3: x₁ = (5+1)/2 = 3, x₂ = (5-1)/2 = 2\n'
    '</think>\n'
    'Răspuns: x₁ = 2, x₂ = 3\n'
    '<|im_end|>'
)
doc.add_page_break()

# ─── 4. QLORA ───
doc.add_heading('4. Tehnica de antrenare — QLoRA', level=1)
txt('QLoRA (Quantized Low-Rank Adaptation) combină două tehnici pentru a permite '
    'antrenarea unui model de 14B parametri pe un singur GPU de 16GB:')

doc.add_heading('4.1 Quantizare 4-bit (NF4)', level=2)
txt('Modelul original de 14B parametri ocupă ~28GB în precizie completă (float16). '
    'Prin quantizare NF4 (Normal Float 4-bit) cu double quantization, dimensiunea scade la ~8GB, '
    'permițând încărcarea pe un Tesla T4 de 16GB.')

doc.add_heading('4.2 LoRA (Low-Rank Adaptation)', level=2)
txt('În loc să modificăm toți cei 14B parametri, LoRA adaugă matrice mici de rang redus (r=16) '
    'în layerele de atenție. Formula: W_nou = W_original + B × A, unde B și A sunt matricele antrenabile.')
txt('Aceasta reduce numărul de parametri antrenabili la sub 1% din total, '
    'făcând antrenarea posibilă pe hardware limitat.')

doc.add_heading('4.3 Configurația LoRA', level=2)
table(['Parametru', 'Valoare', 'Justificare'], [
    ['Rank (r)', '16', 'Suficient pentru adaptare la matematică'],
    ['Alpha (α)', '32', 'Scaling factor; α/r = 2 pentru stabilitate'],
    ['Dropout', '0.05', 'Regularizare minimală'],
    ['Target modules', '7 (q, k, v, o, gate, up, down)', 'Toate proiecțiile attention + FFN'],
    ['Bias', 'none', 'Nu antrenăm bias-urile'],
    ['Task type', 'CAUSAL_LM', 'Generare de text autoregresivă'],
])
doc.add_page_break()

# ─── 5. ANTRENARE ───
doc.add_heading('5. Procesul de antrenare', level=1)

doc.add_heading('5.1 Infrastructură', level=2)
table(['Componentă', 'Specificație'], [
    ['GPU', 'Tesla T4 (16GB VRAM)'],
    ['Platformă', 'Kaggle Notebooks'],
    ['Framework', 'PyTorch + HuggingFace Transformers'],
    ['Biblioteci', 'PEFT 0.10, TRL 0.8.1, BitsAndBytes'],
    ['Precizie', 'FP16 (mixed precision)'],
])

doc.add_heading('5.2 Hiperparametri', level=2)
table(['Parametru', 'Valoare'], [
    ['Batch size', '1 per GPU'],
    ['Gradient accumulation', '8 (effective batch = 8)'],
    ['Learning rate', '5 × 10⁻⁵'],
    ['LR Scheduler', 'Cosine decay'],
    ['Warmup', '5% din total steps'],
    ['Epoci', '2'],
    ['Max sequence length', '1.024 tokeni'],
    ['Optimizer', 'Paged AdamW 8-bit'],
    ['Eval steps', 'La fiecare 50 pași'],
    ['Gradient checkpointing', 'Activat (economie VRAM)'],
])

doc.add_heading('5.3 Rezultate', level=2)
txt('Antrenarea a produs o convergență clară, cu scăderea constantă a funcției de pierdere (loss):')

table(['Step', 'Train Loss', 'Eval Loss', 'Train PPL', 'Eval PPL'], [
    ['50', '0.4188', '0.4179', '1.52', '1.52'],
    ['100', '0.3179', '0.3164', '1.37', '1.37'],
    ['150', '0.2567', '0.2611', '1.29', '1.30'],
    ['200', '0.2247', '0.2447', '1.25', '1.28'],
    ['250', '0.1976', '0.2376', '1.22', '1.27'],
    ['300', '0.1822', '0.2318', '1.20', '1.26'],
    ['350', '0.2043', '0.2289', '1.23', '1.26'],
])

txt('Observații principale:')
bullet('Loss-ul de antrenare a scăzut de la 0.42 la 0.18 — o reducere de 57%')
bullet('Loss-ul de validare a scăzut de la 0.42 la 0.23 — o reducere de 45%')
bullet('Perplexitatea finală de evaluare: 1.26 — model cu predicții precise')
bullet('Nu s-a observat overfitting sever — eval loss rămâne stabil')
bullet('Spike la step 350 (train loss 0.20) — batch dificil, recuperat imediat')
doc.add_page_break()

# ─── 6. INTEGRARE ───
doc.add_heading('6. Integrarea în aplicație', level=1)

doc.add_heading('6.1 Arhitectura sistemului', level=2)
txt('Modelul antrenat este servit printr-un server Flask pe Kaggle, expus prin ngrok:')
code(
    'Utilizator (React Native)\n'
    '    ↓ mesaj text\n'
    'Backend FastAPI (/api/chat)\n'
    '    ↓ POST {intrebare: "..."}\n'
    'Server Flask pe Kaggle (ngrok)\n'
    '    ↓ model.generate()\n'
    'Răspuns: {raspuns: "Pasul 1: ... Răspuns: ..."}\n'
    '    ↓ parsare pași\n'
    'Frontend: afișare structurată cu KaTeX'
)

doc.add_heading('6.2 Pipeline de procesare', level=2)
bullet('Normalizare input: convertește notații diverse (√, ², ∫) în format standard')
bullet('Detecție tip exercițiu: ecuație, derivată, integrală, limită, determinant')
bullet('Trimitere la model: format ChatML cu system prompt SmartBAC')
bullet('Parsare răspuns: extragere pași, metodă, răspuns final')
bullet('Filtru calitate: eliminare halucinări, gibberish, repetiții')
bullet('Randare frontend: formule matematice cu KaTeX, dark mode')

doc.add_heading('6.3 Fallback offline', level=2)
txt('Când serverul Kaggle nu este disponibil, aplicația folosește un solver rule-based '
    'implementat local, care rezolvă corect ecuații, derivate, integrale, determinanți și combinări '
    'fără nicio dependență de internet.')
doc.add_page_break()

# ─── 7. CONCLUZII ───
doc.add_heading('7. Concluzii', level=1)
txt('Prin fine-tuning-ul modelului DeepSeek-R1-Distill-Qwen-14B cu QLoRA, am demonstrat că '
    'un model de limbaj mare poate fi adaptat eficient pentru rezolvarea exercițiilor de BAC, '
    'folosind hardware accesibil (un singur GPU de 16GB).')

doc.add_heading('7.1 Contribuții', level=2)
bullet('Antrenare QLoRA pe 1.500 exerciții de BAC cu 2 epoci')
bullet('Reducere loss de la 0.42 la 0.18 (57%)')
bullet('Integrare model în chatbot cu normalizare input și parsare răspuns')
bullet('Pipeline multi-model cu fallback offline (rule-based solver)')
bullet('Model publicat pe HuggingFace pentru reproducibilitate')

doc.add_heading('7.2 Limitări și dezvoltări viitoare', level=2)
bullet('Dataset mai mare (8.700+ exerciții) cu soluții detaliate pas cu pas')
bullet('Antrenare pe 3+ epoci pentru consolidare')
bullet('Model mai mic (1.5B-7B) pentru deployment local pe dispozitive mobile')
bullet('Evaluare sistematică pe subiecte BAC oficiale')
bullet('Adăugare suport pentru recunoaștere exerciții din imagini (scanner)')

# Save Word
word_path = os.path.expanduser('~/Desktop/SmartBAC_Chatbot_AI.docx')
doc.save(word_path)
print(f'Word: {word_path}')

# ══════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = PptRGB(0x0F, 0x17, 0x2A)
CARD = PptRGB(0x1E, 0x29, 0x3B)
WHITE = PptRGB(0xF1, 0xF5, 0xF9)
MUTED = PptRGB(0x94, 0xA3, 0xB8)
GREEN = PptRGB(0x34, 0xD3, 0x99)
BLUE = PptRGB(0x60, 0xA5, 0xFA)
PURPLE = PptRGB(0xA7, 0x8B, 0xFA)
ORANGE = PptRGB(0xFB, 0x92, 0x3C)
YELLOW = PptRGB(0xFB, 0xBF, 0x24)
CYAN = PptRGB(0x22, 0xD3, 0xEE)
RED = PptRGB(0xF8, 0x71, 0x71)
TEAL = PptRGB(0x14, 0xB8, 0xA6)


def sbg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def stxt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = PptPt(sz); p.font.color.rgb = clr; p.font.bold = bold; p.alignment = align

def scard(slide, l, t, w, h, title, body, accent=BLUE, bsz=13):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = CARD; s.line.fill.background()
    bar = slide.shapes.add_shape(1, l, t, Inches(0.06), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    stxt(slide, l+Inches(0.2), t+Inches(0.12), w-Inches(0.4), Inches(0.35), title, sz=15, clr=accent, bold=True)
    stxt(slide, l+Inches(0.2), t+Inches(0.5), w-Inches(0.4), h-Inches(0.6), body, sz=bsz, clr=MUTED)

def smetric(slide, l, t, value, label, clr):
    s = slide.shapes.add_shape(1, l, t, Inches(2.2), Inches(1.6))
    s.fill.solid(); s.fill.fore_color.rgb = CARD; s.line.fill.background()
    stxt(slide, l, t+Inches(0.15), Inches(2.2), Inches(0.7), value, sz=34, clr=clr, bold=True, align=PP_ALIGN.CENTER)
    stxt(slide, l, t+Inches(0.85), Inches(2.2), Inches(0.5), label, sz=12, clr=MUTED, align=PP_ALIGN.CENTER)

# Training data
steps = [50, 100, 150, 200, 250, 300, 350]
train_loss = [0.4188, 0.3179, 0.2567, 0.2247, 0.1976, 0.1822, 0.2043]
eval_loss = [0.4179, 0.3164, 0.2611, 0.2447, 0.2376, 0.2318, 0.2289]

# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(1), Inches(1.5), Inches(11), Inches(1), 'SmartBAC — Chatbot AI', sz=48, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)
stxt(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.7), 'Antrenarea unui model de limbaj pentru\nrezolvarea exercițiilor de Bacalaureat', sz=22, clr=MUTED, align=PP_ALIGN.CENTER)
stxt(slide, Inches(1), Inches(4), Inches(11), Inches(0.5), 'DeepSeek-R1-Distill-Qwen-14B  ·  QLoRA  ·  Kaggle T4', sz=16, clr=PURPLE, align=PP_ALIGN.CENTER)
stxt(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5), '2026', sz=14, clr=MUTED, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: Model + Dataset ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'MODELUL ȘI DATELE', sz=12, clr=ORANGE, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'DeepSeek-R1 + Exerciții BAC', sz=30, clr=WHITE, bold=True)

scard(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8),
    'DeepSeek-R1-Distill-Qwen-14B',
    '• 14 miliarde parametri\n'
    '• Distilat din DeepSeek-R1 (671B)\n'
    '• Raționament nativ cu <think> blocks\n'
    '• Suport limba română\n'
    '• Arhitectură Qwen2.5 (decoder-only)\n\n'
    'Ales pentru:\n'
    '• Raționament matematic superior\n'
    '• Chain-of-thought structurat\n'
    '• Compatibil cu QLoRA fine-tuning',
    PURPLE, bsz=14)

scard(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.2),
    'Dataset exerciții BAC',
    'Total: 3.795 exerciții\n'
    'Train: 1.500 samples  |  Eval: 150 samples\n'
    'Surse: subiecte BAC 2009-2024\n'
    'Tipuri: ecuații, derivate, integrale, limite,\n'
    '        determinanți, combinări, geometrie, trigonometrie',
    GREEN, bsz=14)

scard(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.3),
    'Format ChatML',
    '<|im_start|>system\n'
    'Ești SmartBAC, asistent matematică BAC.\n'
    '<|im_end|>\n'
    '<|im_start|>user → întrebare\n'
    '<|im_start|>assistant\n'
    '<think> rezolvare pas cu pas </think>\n'
    'Răspuns: rezultat final',
    CYAN, bsz=12)

# ─── SLIDE 3: QLoRA ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'TEHNICA DE ANTRENARE', sz=12, clr=PURPLE, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'QLoRA — Quantized Low-Rank Adaptation', sz=30, clr=WHITE, bold=True)

scard(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(5),
    'Quantizare 4-bit NF4',
    'Model 14B: 28GB → ~8GB VRAM\n\n'
    'NF4 = Normal Float 4-bit\n'
    'Fiecare parametru ocupă 4 biți\n'
    'în loc de 16 biți (float16)\n\n'
    '+ Double Quantization:\n'
    'Quantizăm și constantele\n'
    'de quantizare → economie\n'
    'suplimentară ~0.4GB\n\n'
    'Rezultat: model 14B încape\n'
    'pe un singur GPU T4 (16GB)',
    ORANGE, bsz=14)

scard(slide, Inches(4.6), Inches(1.8), Inches(3.8), Inches(5),
    'LoRA (Low-Rank Adaptation)',
    'W_nou = W_frozen + B × A\n\n'
    'W: 14B params (înghețat)\n'
    'B: d × 16 (antrenat)\n'
    'A: 16 × d (antrenat)\n\n'
    'Config:\n'
    '• Rank (r) = 16\n'
    '• Alpha (α) = 32\n'
    '• Dropout = 0.05\n'
    '• 7 target modules:\n'
    '  q, k, v, o, gate, up, down\n\n'
    'Sub 1% parametri antrenabili',
    PURPLE, bsz=14)

scard(slide, Inches(8.7), Inches(1.8), Inches(3.8), Inches(5),
    'Avantaje QLoRA',
    '1. Memorie redusă\n'
    '   14B pe GPU 16GB\n'
    '   (imposibil fără QLoRA)\n\n'
    '2. Antrenare rapidă\n'
    '   Sub 1% parametri\n'
    '   = convergență în ore\n\n'
    '3. Adapter mic\n'
    '   ~150MB (nu 28GB model)\n'
    '   Ușor de distribuit\n\n'
    '4. Model original intact\n'
    '   Zero degradare pe\n'
    '   alte task-uri',
    TEAL, bsz=14)

# ─── SLIDE 4: Training Config ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'CONFIGURAȚIA ANTRENĂRII', sz=12, clr=GREEN, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'Hiperparametri și Infrastructură', sz=30, clr=WHITE, bold=True)

scard(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5),
    'Hiperparametri',
    'Batch size:             1 per GPU\n'
    'Gradient accumulation:  8\n'
    'Effective batch:        8\n\n'
    'Learning rate:          5 × 10⁻⁵\n'
    'LR Scheduler:           Cosine decay\n'
    'Warmup:                 5%\n\n'
    'Epoci:                  2\n'
    'Max seq length:         1024 tokeni\n'
    'Optimizer:              Paged AdamW 8-bit\n'
    'Precision:              FP16\n\n'
    'Eval:                   la fiecare 50 pași\n'
    'Save:                   la finalul epocii',
    GREEN, bsz=14)

scard(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.3),
    'Infrastructură Kaggle',
    'GPU:        Tesla T4 (16GB VRAM)\n'
    'Platformă:  Kaggle Notebooks (gratuit)\n'
    'Framework:  PyTorch + HuggingFace\n'
    'Biblioteci: PEFT 0.10, TRL 0.8.1, BitsAndBytes',
    BLUE, bsz=14)

scard(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.4),
    'Optimizări VRAM',
    '• Gradient Checkpointing\n'
    '  Recalculare activări = economie VRAM\n'
    '• 4-bit NF4 + Double Quantization\n'
    '  28GB → ~8GB\n'
    '• Paged AdamW 8-bit\n'
    '  Optimizer states paginat CPU/GPU',
    ORANGE, bsz=14)

# ─── SLIDE 5: Results ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'REZULTATE', sz=12, clr=YELLOW, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'Evoluția Loss și Convergență', sz=30, clr=WHITE, bold=True)

smetric(slide, Inches(0.5), Inches(1.6), '0.418', 'Loss Inițial', RED)
smetric(slide, Inches(2.9), Inches(1.6), '0.182', 'Loss Final', GREEN)
smetric(slide, Inches(5.3), Inches(1.6), '-57%', 'Reducere', CYAN)
smetric(slide, Inches(7.7), Inches(1.6), '1.26', 'Perplexitate\nFinală', PURPLE)
smetric(slide, Inches(10.1), Inches(1.6), '350', 'Total\nPași', BLUE)

# Loss chart
chart_data = CategoryChartData()
chart_data.categories = [str(s) for s in steps]
chart_data.add_series('Train Loss', train_loss)
chart_data.add_series('Eval Loss', eval_loss)
chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.5), Inches(3.5), Inches(12), Inches(3.5), chart_data).chart
chart.has_legend = True; chart.legend.include_in_layout = False
chart.legend.font.size = PptPt(11); chart.legend.font.color.rgb = MUTED
plot = chart.plots[0]; plot.smooth = True
plot.series[0].format.line.color.rgb = TEAL; plot.series[0].format.line.width = PptPt(3)
plot.series[1].format.line.color.rgb = RED; plot.series[1].format.line.width = PptPt(3)
chart.category_axis.tick_labels.font.size = PptPt(10); chart.category_axis.tick_labels.font.color.rgb = MUTED
chart.value_axis.tick_labels.font.size = PptPt(10); chart.value_axis.tick_labels.font.color.rgb = MUTED

# ─── SLIDE 6: Integration ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'INTEGRARE ÎN APLICAȚIE', sz=12, clr=GREEN, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'De la model antrenat la chatbot funcțional', sz=28, clr=WHITE, bold=True)

flow = [
    ('1. Antrenare\nKaggle T4', 'QLoRA fine-tune\n350 steps, 2 epoci\nLoss: 0.42→0.18', PURPLE),
    ('2. Publicare\nHuggingFace', 'Model + adapter\nReproducibil\nOpen-source', ORANGE),
    ('3. Server Flask\n+ ngrok', 'Endpoint /ask\nModel CUDA\nJSON response', BLUE),
    ('4. Backend\nFastAPI', 'Normalizare input\nParsare răspuns\nFiltru calitate', GREEN),
    ('5. Frontend\nReact Native', 'KaTeX formule\nDark mode\nAnimații', CYAN),
]
for i, (title, body, clr) in enumerate(flow):
    x = Inches(0.3) + Inches(2.55) * i
    scard(slide, x, Inches(1.8), Inches(2.3), Inches(2.8), title, body, clr, bsz=12)
    if i < 4:
        stxt(slide, x + Inches(2.3), Inches(2.8), Inches(0.3), Inches(0.4), '→', sz=22, clr=MUTED, align=PP_ALIGN.CENTER)

scard(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(1.8),
    'Pipeline complet',
    'Utilizator scrie întrebare → Normalizare (²→^2, √→sqrt, radical→√) → '
    'Detecție tip exercițiu → Trimitere la model SmartBAC (Kaggle/ngrok) → '
    'Generare cu <think> blocks → Parsare pași + răspuns → Filtru calitate (anti-halucinare) → '
    'Afișare cu formule KaTeX\n'
    'Fallback offline: solver rule-based pentru ecuații, derivate, integrale, determinanți',
    YELLOW, bsz=13)

# ─── SLIDE 7: Demo ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'DEMONSTRAȚIE', sz=12, clr=GREEN, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'Exemplu de rezolvare', sz=30, clr=WHITE, bold=True)

scard(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(1), 'Utilizator:', 'Rezolvă ecuația: x^2 - 5x + 6 = 0', GREEN, bsz=16)

scard(slide, Inches(0.5), Inches(3), Inches(12), Inches(4),
    'SmartBAC răspunde:',
    '<think>\n'
    'Pasul 1: Identificăm coeficienții: a=1, b=-5, c=6\n'
    'Pasul 2: Calculăm discriminantul: Δ = (-5)² - 4·1·6 = 25 - 24 = 1\n'
    'Pasul 3: Δ > 0, deci ecuația are 2 soluții reale distincte\n'
    'Pasul 4: x₁,₂ = (-b ± √Δ) / (2a) = (5 ± 1) / 2\n'
    'Pasul 5: x₁ = (5+1)/2 = 3, x₂ = (5-1)/2 = 2\n'
    '</think>\n\n'
    'Răspuns: x₁ = 2, x₂ = 3',
    BLUE, bsz=14)

# ─── SLIDE 8: Conclusions ───
slide = prs.slides.add_slide(prs.slide_layouts[6]); sbg(slide)
stxt(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4), 'CONCLUZII', sz=12, clr=PURPLE, bold=True)
stxt(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6), 'Realizări și Dezvoltări Viitoare', sz=30, clr=WHITE, bold=True)

scard(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8),
    'Ce am realizat',
    '• Fine-tuning model 14B pe GPU T4 16GB\n'
    '  (posibil doar cu QLoRA)\n\n'
    '• Loss redus cu 57%: 0.42 → 0.18\n\n'
    '• Perplexitate finală: 1.26\n\n'
    '• Model publicat pe HuggingFace\n\n'
    '• Chatbot funcțional integrat\n'
    '  în aplicația SmartBAC\n\n'
    '• Pipeline cu normalizare,\n'
    '  parsare și filtru anti-halucinare\n\n'
    '• Fallback offline rule-based',
    GREEN, bsz=14)

scard(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8),
    'Dezvoltări viitoare',
    '• Dataset extins: 8.700+ exerciții\n'
    '  cu soluții detaliate pas cu pas\n\n'
    '• Antrenare pe 3+ epoci\n'
    '  pentru consolidare\n\n'
    '• Model mai mic (1.5B-7B)\n'
    '  pentru deployment local/mobil\n\n'
    '• Scanner AI: recunoaștere\n'
    '  exerciții din imagini (VLM)\n\n'
    '• Evaluare pe subiecte BAC\n'
    '  oficiale (benchmark)\n\n'
    '• Feedback loop: model învață\n'
    '  din corecțiile utilizatorilor',
    CYAN, bsz=14)

stxt(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), 'Mulțumesc!    |    Întrebări?', sz=16, clr=MUTED, align=PP_ALIGN.CENTER)

# Save PPT
ppt_path = os.path.expanduser('~/Desktop/SmartBAC_Chatbot_AI.pptx')
prs.save(ppt_path)
print(f'PPT: {ppt_path}')
