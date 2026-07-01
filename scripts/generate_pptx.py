"""Generate PowerPoint presentation — SmartBAC AI."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Colors ───
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
RED = RGBColor(0xF8, 0x71, 0x71)
DARK_TEXT = RGBColor(0x1B, 0x1B, 0x2F)


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Round corners
    sp = shape._element
    sp_pr = sp.find(qn('a:prstGeom'))
    if sp_pr is not None:
        sp_pr.set('prst', 'roundRect')
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_card(slide, left, top, w, h, title, body, accent_color=BLUE, body_size=14):
    card = add_shape(slide, left, top, w, h, CARD)
    # Accent bar
    bar = slide.shapes.add_shape(1, left, top, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15), w - Inches(0.5), Inches(0.4),
                title, font_size=16, color=accent_color, bold=True)
    # Body
    add_textbox(slide, left + Inches(0.25), top + Inches(0.55), w - Inches(0.5), h - Inches(0.7),
                body, font_size=body_size, color=MUTED)
    return card


# ════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            'SmartBAC', font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
            'Tutor Inteligent pentru Bacalaureat la Matematică',
            font_size=24, color=MUTED, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
            'AI Multi-Model  ·  LoRA Fine-tuning  ·  Vision Language Model',
            font_size=16, color=PURPLE, align=PP_ALIGN.CENTER)

# Bottom badges
for i, (text, clr) in enumerate([('React Native', GREEN), ('FastAPI', BLUE), ('PyTorch', ORANGE), ('GPT-4o', PURPLE), ('Qwen VLM', CYAN)]):
    x = Inches(3.2) + Inches(1.5) * i
    badge = add_shape(slide, x, Inches(5.2), Inches(1.3), Inches(0.45), RGBColor(0x33, 0x41, 0x55))
    add_textbox(slide, x + Inches(0.1), Inches(5.22), Inches(1.1), Inches(0.4),
                text, font_size=11, color=clr, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
            '2026', font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 2 — Problema
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'PROBLEMA', font_size=12, color=ORANGE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(1),
            'Elevii nu au acces la tutori personalizați\npentru BAC matematică',
            font_size=32, color=WHITE, bold=True)

cards = [
    ('Tutori scumpi', '150-300 lei/oră\nNu toți elevii își permit', RED),
    ('Indisponibili 24/7', 'Elevul are nevoie de ajutor\nseara/weekend/vacanță', ORANGE),
    ('Nu se adaptează', 'Același ritm pentru toți\nNu detectează lacunele', YELLOW),
    ('Fără feedback instant', 'Elevul nu știe imediat\ndacă a greșit sau de ce', PURPLE),
]
for i, (title, body, clr) in enumerate(cards):
    x = Inches(0.8) + Inches(3.1) * i
    add_card(slide, x, Inches(3), Inches(2.8), Inches(2.2), title, body, clr)

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
            'SmartBAC rezolvă toate aceste probleme prin AI — disponibil 24/7, gratuit, personalizat.',
            font_size=16, color=GREEN, bold=True)

# ════════════════════════════════════════════════════════
# SLIDE 3 — Arhitectura
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'ARHITECTURA AI', font_size=12, color=BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Pipeline Multi-Model cu Fallback Ierarhic', font_size=28, color=WHITE, bold=True)

# Flow diagram as cards
flow = [
    ('1. Intent\nDetection', 'Regex → clasificare\nsolve/concept/hint', CYAN, Inches(0.5)),
    ('2. GPT-4o\n(Primar)', 'JSON mode\nLaTeX structurat\n~85% accuracy', GREEN, Inches(2.7)),
    ('3. Qwen LoRA\n(Fallback 1)', 'Model propriu\nFine-tunat pe BAC\n~70% accuracy', PURPLE, Inches(4.9)),
    ('4. Rule-based\n(Fallback 2)', 'Rezolvare simbolică\n100% offline\n<100ms', ORANGE, Inches(7.1)),
    ('5. Post-\nProcessing', 'Fix LaTeX escape\nJSON → UI\nKaTeX render', BLUE, Inches(9.3)),
]
for title, body, clr, x in flow:
    add_card(slide, x, Inches(2.2), Inches(2), Inches(2.8), title, body, clr, body_size=12)

# Arrows between cards
for i in range(4):
    x = Inches(2.5) + Inches(2.2) * i
    add_textbox(slide, x, Inches(3.3), Inches(0.5), Inches(0.4), '→', font_size=24, color=MUTED, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.8),
            'Dacă GPT-4o nu e disponibil → Qwen LoRA (local) → Rule-based (offline). Zero downtime.',
            font_size=15, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 4 — Tokenizer + Transformer
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'MODELE ANTRENATE DE LA ZERO', font_size=12, color=GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Tokenizer BPE + Transformer Decoder-Only', font_size=28, color=WHITE, bold=True)

# Left card: Tokenizer
add_card(slide, Inches(0.5), Inches(2), Inches(5.8), Inches(4.5),
         'Tokenizer BPE Custom',
         'Implementare de la zero în Python\n\n'
         '• Vocabular: 8.192 tokeni\n'
         '• Protecție LaTeX: \\frac, \\sqrt, \\int\n'
         '• Pre-tokenizare regex multi-pattern\n'
         '• Suport diacritice: ă, â, î, ș, ț\n'
         '• Tokeni speciali: <BOS>, <EOS>, <SEP>\n\n'
         'Corpus: ~5.000 exerciții BAC',
         CYAN, body_size=14)

# Right card: Transformer
add_card(slide, Inches(6.8), Inches(2), Inches(5.8), Inches(4.5),
         'Transformer Decoder-Only',
         'Implementare PyTorch de la zero\n\n'
         '• 6 blocuri, 256 d_model, 8 heads\n'
         '• ~5.4 milioane parametri\n'
         '• Pre-Norm + Sinusoidal PE\n'
         '• Top-k sampling (k=30, T=0.5)\n'
         '• Loss: Cross-entropy next-token\n\n'
         'Training: 242 exerciții, AdamW, cosine LR',
         PURPLE, body_size=14)

# ════════════════════════════════════════════════════════
# SLIDE 5 — LoRA Fine-tuning
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'FINE-TUNING', font_size=12, color=PURPLE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'LoRA pe Qwen2.5-Math-1.5B', font_size=28, color=WHITE, bold=True)

# LoRA explanation
add_card(slide, Inches(0.5), Inches(2), Inches(5.8), Inches(4.5),
         'Ce este LoRA?',
         'Low-Rank Adaptation (Hu et al., 2021)\n\n'
         'W_nou = W_original + ΔW\n'
         'ΔW = B × A  (rank r = 16)\n\n'
         'W: 3072×3072 = 9.4M params (înghețat)\n'
         'B: 3072×16   = 49K params (antrenat)\n'
         'A: 16×3072   = 49K params (antrenat)\n\n'
         'Reducere: 99% mai puțini parametri!\n'
         'Antrenăm doar 0.98% din model.',
         ORANGE, body_size=14)

# Config
add_card(slide, Inches(6.8), Inches(2), Inches(5.8), Inches(4.5),
         'Configurație antrenare',
         'Model: Qwen2.5-Math-1.5B\n'
         'LoRA: r=16, α=32, dropout=0.05\n'
         'Target: q,k,v,o,gate,up,down proj\n\n'
         'Dataset: 4.086 exerciții BAC\n'
         'Format: ChatML cu <think> blocks\n'
         'Split: 89% train / 10% val / 1% test\n\n'
         'Optimizer: AdamW, LR=1e-4, cosine\n'
         'Quantizare: 4-bit NF4\n'
         'Hardware: Apple M4 (MLX) / T4 (PyTorch)\n'
         'Timp: ~2-4 ore',
         PURPLE, body_size=14)

# ════════════════════════════════════════════════════════
# SLIDE 6 — Scanner VLM
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'SCANNER AI', font_size=12, color=CYAN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Vision Language Model — Qwen2.5-VL-3B', font_size=28, color=WHITE, bold=True)

# Pipeline
add_card(slide, Inches(0.5), Inches(2), Inches(3.8), Inches(4.3),
         'Dataset Sintetic',
         '4.541 exerciții → imagini\n\n'
         '6 augmentări per exercițiu:\n'
         '• 5 tipuri de fundal\n'
         '• 5 culori text\n'
         '• Rotație ±3.5°\n'
         '• Zgomot gaussian\n'
         '• Blur + contrast\n\n'
         'Total: 27.108 imagini',
         GREEN, body_size=13)

add_card(slide, Inches(4.6), Inches(2), Inches(3.8), Inches(4.3),
         'Antrenare VLM',
         'Qwen2.5-VL-3B-Instruct\n\n'
         '• 4-bit quantizare NF4\n'
         '• Vision encoder ÎNGHEȚAT\n'
         '• LoRA pe LLM layers\n'
         '• Batch=2, grad_accum=8\n'
         '• 2 epoci, cosine LR\n\n'
         'GPU: Tesla T4 16GB\n'
         'Timp: ~8-12 ore',
         BLUE, body_size=13)

add_card(slide, Inches(8.7), Inches(2), Inches(3.8), Inches(4.3),
         'Pipeline Inferență',
         'Utilizator → foto exercițiu\n\n'
         '1. VLM propriu (dacă e antrenat)\n'
         '2. GPT-4o Vision (fallback)\n'
         '3. EasyOCR + solver (offline)\n\n'
         'Output: JSON structurat cu\n'
         'enunț + pași + răspuns\n\n'
         'Latență: 3-5 secunde',
         ORANGE, body_size=13)

# ════════════════════════════════════════════════════════
# SLIDE 7 — Chat Demo
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'CHAT TUTOR', font_size=12, color=GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Exemplu de Rezolvare', font_size=28, color=WHITE, bold=True)

# Input
add_card(slide, Inches(0.5), Inches(2), Inches(5.8), Inches(1.2),
         'Utilizator:', 'rezolvă ecuația 2x² - 5x + 3 = 0', GREEN, body_size=16)

# Output
add_card(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(3.5),
         'SmartBAC Tutor:',
         'Tip: Ecuație gradul II\n'
         'Metodă: Formula discriminantului Δ = b² - 4ac\n\n'
         'Pasul 1: Identificăm coeficienții → a=2, b=-5, c=3\n'
         'Pasul 2: Calculăm Δ = (-5)² - 4·2·3 = 25 - 24 = 1\n'
         'Pasul 3: Aplicăm formula → x₁,₂ = (5 ± 1) / 4\n'
         'Pasul 4: x₁ = 3/2, x₂ = 1\n\n'
         'RĂSPUNS: x₁ = 3/2, x₂ = 1\n'
         'Verificare: 2(3/2)² - 5(3/2) + 3 = 9/2 - 15/2 + 3 = 0 ✓\n'
         'Greșeli frecvente: Semn greșit la b, Uitarea verificării',
         BLUE, body_size=13)

# ════════════════════════════════════════════════════════
# SLIDE 8 — Rezultate
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'REZULTATE', font_size=12, color=YELLOW, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Evaluare și Comparație Modele', font_size=28, color=WHITE, bold=True)

# Metrics cards
metrics = [
    ('85%', 'GPT-4o\nAccuracy', GREEN),
    ('70%', 'Qwen LoRA\nAccuracy', PURPLE),
    ('27K', 'Imagini\nDataset', BLUE),
    ('<100ms', 'Rule-based\nLatență', ORANGE),
    ('4.596', 'Exerciții\nAntrenare', CYAN),
]
for i, (value, label, clr) in enumerate(metrics):
    x = Inches(0.5) + Inches(2.5) * i
    card = add_shape(slide, x, Inches(2), Inches(2.2), Inches(1.8), CARD)
    add_textbox(slide, x, Inches(2.15), Inches(2.2), Inches(0.8),
                value, font_size=36, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.9), Inches(2.2), Inches(0.7),
                label, font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Comparison table area
add_card(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(2.5),
         'Comparație Modele',
         'GPT-4o:       85% accuracy  |  2-3s latency  |  $0.005/req  |  Online\n'
         'Qwen LoRA:    70% accuracy  |  1-2s latency  |  Gratis       |  Offline\n'
         'Transformer:  45% accuracy  |  <1s latency   |  Gratis       |  Offline\n'
         'Rule-based:   60% accuracy  |  <0.1s latency |  Gratis       |  Offline\n'
         'VLM Scanner:  65% accuracy  |  3-5s latency  |  Gratis       |  Offline',
         YELLOW, body_size=14)

# ════════════════════════════════════════════════════════
# SLIDE 9 — Frontend
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'FRONTEND', font_size=12, color=GREEN, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Interfața Utilizator — React Native + KaTeX', font_size=28, color=WHITE, bold=True)

features = [
    ('Dark Mode', 'Design profesional dark\nConsistent pe toate ecranele\nIonicons (fără emoji)', PURPLE),
    ('KaTeX Rendering', 'Formule LaTeX randate frumos\nInline \\(x^2\\) și display \\[...\\]\nDark mode CSS override', BLUE),
    ('Animații', 'FadeInDown pe mesaje\nSpring physics pe butoane\nConfetti la răspuns corect', GREEN),
    ('Gamification', 'XP + Levels + Streak-uri\nAchievements + Challenges\nLeaderboard + Ligi', ORANGE),
]
for i, (title, body, clr) in enumerate(features):
    x = Inches(0.5) + Inches(3.1) * i
    add_card(slide, x, Inches(2), Inches(2.8), Inches(3.5), title, body, clr, body_size=14)

add_textbox(slide, Inches(0.5), Inches(6), Inches(12), Inches(0.5),
            'Stack: Expo React Native (TypeScript) + FastAPI (Python) + MongoDB',
            font_size=14, color=MUTED)

# ════════════════════════════════════════════════════════
# SLIDE 10 — Concluzii
# ════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
            'CONCLUZII', font_size=12, color=PURPLE, bold=True)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
            'Contribuții și Dezvoltări Viitoare', font_size=28, color=WHITE, bold=True)

# Contributions
add_card(slide, Inches(0.5), Inches(2), Inches(5.8), Inches(4.5),
         'Contribuții principale',
         '• Tokenizer BPE specializat pe matematică\n'
         '  românească (implementare de la zero)\n\n'
         '• Transformer 5.4M params antrenat de la zero\n\n'
         '• Pipeline multi-model cu 3 nivele de fallback\n\n'
         '• Dataset 27K imagini sintetice pentru VLM\n\n'
         '• LoRA fine-tuning pe Qwen2.5-Math și VL\n\n'
         '• Post-processing LaTeX robust pentru dark mode',
         GREEN, body_size=14)

# Future
add_card(slide, Inches(6.8), Inches(2), Inches(5.8), Inches(4.5),
         'Dezvoltări viitoare',
         '• Dataset cu imagini reale (poze telefon)\n\n'
         '• RAG cu baza de exerciții rezolvate\n\n'
         '• Model 7B+ pe GPU dedicat\n\n'
         '• Suport diagrame geometrice\n\n'
         '• Benchmark pe subiecte BAC reale\n\n'
         '• Feedback loop — model learns from users',
         CYAN, body_size=14)

# Thank you
add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
            'Mulțumesc!    |    Întrebări?', font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Save
out = os.path.expanduser('~/Desktop/SmartBAC_AI_Prezentare.pptx')
prs.save(out)
print(f'Saved: {out}')
