"""Generează prezentarea SmartBAC (.pptx) — 6 minute, structura facultății."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DIAG = os.path.join(ROOT, "docs", "diagrams", "png")

# Paleta — contrast bun pe videoproiector
NAVY = RGBColor(0x1A, 0x23, 0x7E)   # titluri
DARK = RGBColor(0x20, 0x20, 0x20)   # text
ACCENT = RGBColor(0x2E, 0x7D, 0x32)  # verde accent
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=WHITE):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT,
         bullet=False, space_after=8, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def title_bar(s, text):
    # bandă titlu sus
    tf = box(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    para(tf, text, size=30, color=NAVY, bold=True, first=True)
    # linie sub titlu
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(12.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()


def node(s, l, t, w, h, fill, border, title, sub=None):
    """Casetă rotunjită în stilul deck-ului, cu titlu + subtext."""
    from pptx.util import Inches as I
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(l), I(t), I(w), I(h))
    shp.adjustments[0] = 0.10
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(2.5)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(16)
    r.font.color.rgb = border; r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(3)
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK; r2.font.name = "Calibri"
    return shp


def arrow(s, x1, y1, x2, y2, color=GRAY, label=None):
    """Săgeată dreaptă cu vârf, opțional cu etichetă la mijloc."""
    from pptx.util import Inches as I
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(2.5)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'lg', 'len': 'lg'}))
    if label:
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        tb = box(s, I(mx - 1.0), I(my - 0.32), I(2.0), I(0.4))
        tb.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = tb.paragraphs[0].add_run(); rr.text = label
        rr.font.size = Pt(11); rr.font.italic = True; rr.font.color.rgb = GRAY; rr.font.name = "Calibri"


def add_image_fit(s, path, l, t, max_w, max_h):
    from PIL import Image
    if not os.path.exists(path):
        return
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    left = l + (max_w - w) // 2
    top = t + (max_h - h) // 2
    s.shapes.add_picture(path, Emu(left), Emu(top), Emu(w), Emu(h))


# ---------- SLIDE 1: TITLU ----------
s = slide(); bg(s)
tf = box(s, Inches(1.0), Inches(0.5), Inches(11.3), Inches(1.6))
para(tf, "Universitatea Transilvania din Brașov", size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "Facultatea de Inginerie Electrică și Știința Calculatoarelor", size=15, color=GRAY, align=PP_ALIGN.CENTER)
para(tf, "Departamentul Automatică și Tehnologia Informației  ·  Programul de studii: Tehnologia Informației", size=13, color=GRAY, align=PP_ALIGN.CENTER)

tf = box(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.2))
para(tf, "SmartBAC", size=54, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "Aplicație mobilă bazată pe inteligență artificială pentru\npregătirea examenului de Bacalaureat la matematică", size=22, color=DARK, bold=True, align=PP_ALIGN.CENTER)

tf = box(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.5))
para(tf, "Absolvent: Stan Ioana Denisa", size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "Conducător științific: Prof. dr. ing. Kristály Dominic Mircea", size=17, color=GRAY, align=PP_ALIGN.CENTER)
para(tf, "Brașov, 2026", size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ---------- SLIDE 2: TEMA ----------
s = slide(); bg(s); title_bar(s, "Tema proiectului")
tf = box(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
para(tf, "Aplicație mobilă pentru pregătirea BAC-ului la matematică (M1/M2)", size=22, bold=True, color=DARK, bullet=True, first=True, space_after=16)
para(tf, "Rezolvă exerciții pas cu pas — 4.541 exerciții BAC (2008–2025)", size=20, bullet=True, space_after=12)
para(tf, "Explică prin chat cu un model AI antrenat de mine", size=20, bullet=True, space_after=12)
para(tf, "Scanează exerciții din poze (OCR) cu un model AI antrenat de mine", size=20, bullet=True, space_after=12)
para(tf, "Prezice nota la BAC pe baza performanței elevului", size=20, bullet=True, space_after=12)
para(tf, "Motivează prin gamificare: XP, nivele, ligi, streak-uri", size=20, bullet=True, space_after=12)
para(tf, "Problema rezolvată: feedback personalizat + motivație constantă, într-o singură aplicație", size=20, bold=True, color=ACCENT, space_after=8)

# ---------- SLIDE 3: TEHNOLOGII ----------
s = slide(); bg(s); title_bar(s, "Tehnologii folosite și motivația alegerii")
rows = [
    ("Frontend", "React Native + Expo", "un singur cod pentru Android și iOS (în loc de 2 aplicații native separate)"),
    ("Backend", "FastAPI", "API REST asincron + comunicare ușoară cu modulele AI (mai rapid decât Flask)"),
    ("Bază de date", "MongoDB", "model NoSQL flexibil pentru exerciții cu structură variabilă (față de tabelele rigide SQL)"),
    ("Inteligență artificială", "PyTorch, Transformers, QLoRA", "antrenarea și integrarea modelelor mari cu resurse hardware reduse (GPU gratuit de 16 GB)"),
    ("Predicția notei", "Scikit-learn (regresie)", "estimarea notei pe baza progresului și rezultatelor elevului"),
]
from pptx.util import Inches as I
tbl = s.shapes.add_table(len(rows) + 1, 3, I(0.6), I(1.7), I(12.1), I(5.0)).table
tbl.columns[0].width = I(2.7); tbl.columns[1].width = I(3.6); tbl.columns[2].width = I(5.8)
hdr = ["Componentă", "Tehnologie utilizată", "Motivația alegerii"]
for j, h in enumerate(hdr):
    c = tbl.cell(0, j); c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    p = c.text_frame.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE
for i, r in enumerate(rows, 1):
    for j, val in enumerate(r):
        c = tbl.cell(i, j); c.text = val
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        p = c.text_frame.paragraphs[0]; p.font.size = Pt(14)
        p.font.bold = (j == 0); p.font.color.rgb = DARK if j else NAVY

# ---------- SLIDE 4: ARHITECTURA (desenată nativ) ----------
BLUE_F, BLUE_B = RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0x15, 0x65, 0xC0)
GREEN_F, GREEN_B = RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x2E, 0x7D, 0x32)
ORANGE_F, ORANGE_B = RGBColor(0xFF, 0xF3, 0xE0), RGBColor(0xE6, 0x51, 0x00)
PURPLE_F, PURPLE_B = RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x6A, 0x1B, 0x9A)

s = slide(); bg(s); title_bar(s, "Arhitectura sistemului")
# Cele 4 zone
node(s, 0.5, 3.0, 2.9, 1.6, BLUE_F, BLUE_B,
     "Aplicație mobilă", "React Native (Expo)\n16 ecrane")
node(s, 4.6, 2.4, 3.6, 2.7, GREEN_F, GREEN_B,
     "Backend — FastAPI", "auth · exerciții · gamificare\npredicție · chat · scanare")
node(s, 9.6, 1.7, 3.2, 1.7, ORANGE_F, ORANGE_B,
     "MongoDB", "users · exercises · attempts\n+ 5 colecții")
node(s, 9.6, 4.6, 3.2, 1.8, PURPLE_F, PURPLE_B,
     "Modele AI — Kaggle GPU", "DeepSeek-R1 · Qwen3-VL\n(Flask + ngrok)")
# Săgeți
arrow(s, 3.4, 3.8, 4.6, 3.8, BLUE_B, "REST / JSON")
arrow(s, 8.2, 3.2, 9.6, 2.55, ORANGE_B, "citește / scrie")
arrow(s, 8.2, 4.3, 9.6, 5.2, PURPLE_B, "chat · scanare")
# Subtitlu flux
tf = box(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6))
para(tf, "Flux de date:  acțiune elev → aplicație → FastAPI → serviciu → MongoDB / AI → răspuns",
     size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER, first=True)

# ---------- SLIDE 5: AI ----------
s = slide(); bg(s); title_bar(s, "Componenta AI — 3 modele antrenate de mine")
add_image_fit(s, os.path.join(DIAG, "Componente AI Pipeline.png"),
              int(I(0.6)), int(I(1.6)), int(I(7.0)), int(I(5.6)))
tf = box(s, Inches(7.9), Inches(1.7), Inches(5.0), Inches(5.3))
para(tf, "1. Rezolvare / Chat", size=17, bold=True, color=NAVY, first=True, space_after=2)
para(tf, "DeepSeek-R1-Distill-Qwen-14B + QLoRA, raționament <think>", size=14, bullet=True, space_after=12)
para(tf, "2. Scanare poze (OCR)", size=17, bold=True, color=NAVY, space_after=2)
para(tf, "Qwen3-VL-8B + QLoRA — extrage exercițiul din imagine", size=14, bullet=True, space_after=12)
para(tf, "3. Predicția notei", size=17, bold=True, color=NAVY, space_after=2)
para(tf, "Ensemble ML (11 modele, scikit-learn) — R² = 0.97", size=14, bullet=True, space_after=12)
para(tf, "Mecanism de siguranță:", size=16, bold=True, color=ACCENT, space_after=2)
para(tf, "AI → MathTutor (reguli) → mesaj generic. Aplicația funcționează mereu.", size=14, bullet=True)

# ---------- SLIDE 6: PREDICTIE + DB ----------
s = slide(); bg(s); title_bar(s, "Predicția notei și schema bazei de date")
tf = box(s, Inches(0.8), Inches(1.7), Inches(6.0), Inches(5.2))
para(tf, "Predicția notei (Ensemble ML):", size=19, bold=True, color=NAVY, first=True, space_after=10)
para(tf, "30 de caracteristici din performanța elevului", size=17, bullet=True, space_after=8)
para(tf, "11 modele candidate → Top-5 selectate automat", size=17, bullet=True, space_after=8)
para(tf, "Combinare prin Voting / Stacking", size=17, bullet=True, space_after=8)
para(tf, "Rezultat: R² = 0.974, eroare medie ~0.4 puncte", size=17, bold=True, color=ACCENT, bullet=True)
add_image_fit(s, os.path.join(DIAG, "Schema Baza de Date SmartBAC.png"),
              int(I(7.0)), int(I(1.7)), int(I(6.0)), int(I(5.3)))

# ---------- SLIDE 7: DETALII DEZVOLTARE ----------
s = slide(); bg(s); title_bar(s, "Alte detalii de proiectare")
tf = box(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
para(tf, "Învățare adaptivă: algoritmul SM-2 (repetiție spațiată) — exercițiile greșite revin mai des", size=20, bullet=True, first=True, space_after=14)
para(tf, "Gamificare: XP, 10 nivele, 8 realizări, ligi competitive săptămânale, streak-uri", size=20, bullet=True, space_after=14)
para(tf, "Antrenare AI: Kaggle GPU T4, cuantizare 4-biți (NF4) pentru a încăpea modelul de 14B", size=20, bullet=True, space_after=14)
para(tf, "Librării: PyTorch, Transformers, PEFT, TRL, scikit-learn, FastAPI, Expo", size=20, bullet=True, space_after=14)
para(tf, "Securitate: autentificare JWT, parole criptate (bcrypt)", size=20, bullet=True)

# ---------- SLIDE 8: DEMONSTRATIE ----------
s = slide(); bg(s, NAVY)
tf = box(s, Inches(1.0), Inches(2.8), Inches(11.3), Inches(2.0))
para(tf, "Demonstrație", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "(aplicația live)", size=24, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# ---------- SLIDE 9: CONCLUZII ----------
s = slide(); bg(s); title_bar(s, "Concluzii")
tf = box(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
para(tf, "Aplicație completă: rezolvare AI + predicția notei + gamificare", size=21, bullet=True, first=True, space_after=14)
para(tf, "Două modele AI fine-tunate de la zero pe date românești de BAC", size=21, bullet=True, space_after=14)
para(tf, "Arhitectură fiabilă cu mecanism de fallback (funcționează și fără AI)", size=21, bullet=True, space_after=18)
para(tf, "Direcții viitoare:", size=21, bold=True, color=NAVY, space_after=6)
para(tf, "extindere la alte materii · mod offline · generare automată de exerciții", size=19, bullet=True, space_after=14)
para(tf, "Valorificare: publicare în App Store / Google Play", size=21, bold=True, color=ACCENT, bullet=True)

# ---------- SLIDE 10: BIBLIOGRAFIE ----------
s = slide(); bg(s); title_bar(s, "Bibliografie selectivă")
tf = box(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
refs = [
    "DeepSeek-AI — DeepSeek-R1: Incentivizing Reasoning Capability in LLMs (2025)",
    "Dettmers et al. — QLoRA: Efficient Finetuning of Quantized LLMs (2023)",
    "Hu et al. — LoRA: Low-Rank Adaptation of Large Language Models (2021)",
    "Wozniak — SuperMemo 2: algoritmul de repetiție spațiată (SM-2)",
    "Documentația oficială React Native, Expo, FastAPI, scikit-learn",
]
first = True
for r in refs:
    para(tf, r, size=18, color=DARK, bullet=True, first=first, space_after=14)
    first = False

# ---------- SLIDE 11: MULTUMESC ----------
s = slide(); bg(s, NAVY)
tf = box(s, Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.5))
para(tf, "Vă mulțumesc pentru timpul acordat.", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)

out = os.path.join(ROOT, "docs", "Prezentare_SmartBAC.pptx")
prs.save(out)
print("✅ Salvat:", out)
print("   Slide-uri:", len(prs.slides._sldIdLst))
